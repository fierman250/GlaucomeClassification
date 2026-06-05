import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pptx import Presentation

prs = Presentation(r"D:\OneDrive\A4_JobREF\9Z_CodeWork\2403CW_GlaucomeClassification\Python_Code\Glaucoma_Classification_Presentation.pptx")

slide16 = prs.slides[15]
print("=== Slide 16 shapes ===")
for shape in slide16.shapes:
    print(f"  type={shape.shape_type}  name={shape.name}")
    if shape.shape_type == 19:
        tbl = shape.table
        nrows = len(tbl.rows)
        ncols = len(tbl.columns)
        print(f"    => Table {nrows} rows x {ncols} cols")
        for ri in range(min(3, nrows)):
            cells = [tbl.rows[ri].cells[ci].text[:12] for ci in range(ncols)]
            print(f"    row {ri}: {cells}")

slide13 = prs.slides[12]
print("\n=== Slide 13 text boxes ===")
for shape in slide13.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t:
            print(f"  [{len(t)} chars] {t[:120]}")
